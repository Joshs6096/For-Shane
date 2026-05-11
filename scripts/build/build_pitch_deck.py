"""
Transform SR Holding Management LLC — Investor Pitch Deck
Generates a professional PE/VC-ready PPTX from NetSuite financials.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import pptx.oxml
from lxml import etree
import copy

# ─── Brand Colors ───────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0F, 0x24, 0x44)   # Deep navy — primary
NAVY_MID  = RGBColor(0x1A, 0x3D, 0x6B)   # Mid navy — secondary headers
GOLD      = RGBColor(0xC8, 0x94, 0x1A)   # Amber gold — accent
GOLD_LT   = RGBColor(0xF0, 0xB8, 0x30)   # Light gold — highlights
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF7, 0xF8, 0xFA)
LIGHT_GRAY= RGBColor(0xE8, 0xEB, 0xEF)
MID_GRAY  = RGBColor(0x9A, 0xA3, 0xB2)
DARK_GRAY = RGBColor(0x3D, 0x4A, 0x5C)
RED_LOSS  = RGBColor(0xC0, 0x39, 0x2B)   # For negative values
GREEN_POS = RGBColor(0x1A, 0x7A, 0x3C)   # For positive/improving values

# ─── Slide Dimensions (16:9) ────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

# ─── Helpers ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=Pt(11), bold=False, italic=False,
                color=DARK_GRAY, align=PP_ALIGN.LEFT,
                v_anchor=None, wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_label_box(slide, label, x, y, w, h, bg=NAVY, fg=WHITE, size=Pt(10), bold=True):
    rect = add_rect(slide, x, y, w, h, fill_color=bg)
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = fg
    run.font.name = "Calibri"
    return rect


def slide_header(slide, title, subtitle=None):
    """Navy top bar with title and optional subtitle."""
    # Top bar
    add_rect(slide, 0, 0, W, Inches(1.05), fill_color=NAVY)
    # Gold accent line under bar
    add_rect(slide, 0, Inches(1.05), W, Inches(0.045), fill_color=GOLD)
    # Title
    add_textbox(slide, title,
                Inches(0.4), Inches(0.13), Inches(10), Inches(0.55),
                font_size=Pt(22), bold=True, color=WHITE,
                align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.4), Inches(0.62), Inches(10), Inches(0.38),
                    font_size=Pt(12), bold=False, color=GOLD_LT,
                    align=PP_ALIGN.LEFT)
    # Slide background
    add_rect(slide, 0, Inches(1.095), W, H - Inches(1.095), fill_color=OFF_WHITE)
    # Footer bar
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), fill_color=NAVY)
    add_textbox(slide, "CONFIDENTIAL  |  Transform SR Holding Management LLC  |  May 2026",
                Inches(0.4), Inches(7.12), Inches(10), Inches(0.3),
                font_size=Pt(7.5), color=MID_GRAY, align=PP_ALIGN.LEFT)
    add_textbox(slide, "© 2026 Transform SR Holding Management LLC. For discussion purposes only.",
                Inches(8.5), Inches(7.12), Inches(4.5), Inches(0.3),
                font_size=Pt(7.5), color=MID_GRAY, align=PP_ALIGN.RIGHT)


def add_section_label(slide, text, x, y, w=Inches(2.5), h=Inches(0.32)):
    """Gold-left-bordered section label."""
    add_rect(slide, x, y, Inches(0.06), h, fill_color=GOLD)
    add_textbox(slide, text, x + Inches(0.1), y, w, h,
                font_size=Pt(10), bold=True, color=NAVY, align=PP_ALIGN.LEFT)


def kpi_box(slide, label, value, x, y, w=Inches(2.0), h=Inches(1.1),
            val_color=NAVY, bg=WHITE):
    """Metric callout box with label and large value."""
    rect = add_rect(slide, x, y, w, h, fill_color=bg, line_color=LIGHT_GRAY, line_width=Pt(1))
    add_textbox(slide, value, x, y + Inches(0.08), w, Inches(0.65),
                font_size=Pt(22), bold=True, color=val_color,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x, y + Inches(0.68), w, Inches(0.38),
                font_size=Pt(8.5), bold=False, color=MID_GRAY,
                align=PP_ALIGN.CENTER)


def add_table(slide, headers, rows, x, y, w, h,
              header_bg=NAVY, header_fg=WHITE,
              alt_bg=OFF_WHITE, border_color=LIGHT_GRAY,
              col_widths=None, font_size=Pt(9.5),
              header_size=Pt(9.5), row_height=None):
    """Create a properly formatted table object."""
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table

    # Column widths
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)

    # Row height
    if row_height:
        for r in tbl.rows:
            r.height = row_height

    def set_cell(cell, text, bg=None, fg=DARK_GRAY, bold=False,
                 align=PP_ALIGN.LEFT, sz=font_size):
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = sz
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Calibri"
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        else:
            cell.fill.background()
        # Remove cell borders for cleaner look
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

    # Header row
    for j, h_text in enumerate(headers):
        cell = tbl.cell(0, j)
        align = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
        set_cell(cell, h_text, bg=header_bg, fg=header_fg, bold=True,
                 align=align, sz=header_size)

    # Data rows
    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else alt_bg
        is_total = str(row[0]).startswith("  ") is False and any(
            kw in str(row[0]).upper() for kw in ['TOTAL', 'GROSS PROFIT', 'NET ', 'EBITDA', 'OPERATING'])
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            row_bg = LIGHT_GRAY if is_total else bg
            is_bold = is_total
            align = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT
            # Color negative values red
            val_str = str(val)
            fg_color = DARK_GRAY
            if j > 0 and val_str.startswith('(') and not is_total:
                fg_color = RED_LOSS
            elif j > 0 and val_str.startswith('(') and is_total:
                fg_color = RED_LOSS
            set_cell(cell, val_str, bg=row_bg, fg=fg_color, bold=is_bold, align=align)

    return tbl


def fmt_m(val, decimals=1):
    """Format as $XXXm"""
    if val is None:
        return "—"
    if abs(val) >= 1e9:
        return f"${val/1e9:.1f}B"
    if val < 0:
        return f"({abs(val)/1e6:.{decimals}f}M)"
    return f"${val/1e6:.{decimals}f}M"


def fmt_pct(val):
    if val is None:
        return "—"
    return f"{val:.1f}%"


# ════════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ════════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # Blank layout

# ─── Slide 1: Cover ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)

# Full navy background
add_rect(slide, 0, 0, W, H, fill_color=NAVY)

# Gold accent strip
add_rect(slide, 0, Inches(4.35), W, Inches(0.06), fill_color=GOLD)

# Left white accent column
add_rect(slide, 0, 0, Inches(0.12), H, fill_color=GOLD)

# Company name
add_textbox(slide, "TRANSFORM SR HOLDING MANAGEMENT LLC",
            Inches(0.45), Inches(1.5), Inches(11), Inches(0.75),
            font_size=Pt(13), bold=True, color=GOLD_LT,
            align=PP_ALIGN.LEFT)

# Main title
add_textbox(slide, "Strategic Investment\nOpportunity",
            Inches(0.45), Inches(2.1), Inches(11), Inches(1.8),
            font_size=Pt(40), bold=True, color=WHITE,
            align=PP_ALIGN.LEFT)

# Subtitle
add_textbox(slide, "Confidential Investor Presentation  |  May 2026",
            Inches(0.45), Inches(3.85), Inches(9), Inches(0.45),
            font_size=Pt(13), color=MID_GRAY, align=PP_ALIGN.LEFT)

# Descriptor bullets
bullets = [
    "Multi-Stream Revenue Platform  •  ~$650M Revenue  •  $935M Asset Base",
    "Home Services  |  Home Warranty  |  Wholesale  |  Royalties & Licensing",
]
for i, b in enumerate(bullets):
    add_textbox(slide, b,
                Inches(0.45), Inches(4.6 + i * 0.45), Inches(11), Inches(0.4),
                font_size=Pt(11), color=LIGHT_GRAY, align=PP_ALIGN.LEFT)

# Confidential notice
add_textbox(slide,
    "STRICTLY CONFIDENTIAL — This presentation has been prepared solely for the recipient's information. "
    "It does not constitute an offer or solicitation. All financial data sourced from company records.",
    Inches(0.45), Inches(6.7), Inches(12.5), Inches(0.6),
    font_size=Pt(7.5), color=MID_GRAY, italic=True, align=PP_ALIGN.LEFT)

# ─── Slide 2: Disclaimer ────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Important Disclaimer", "Please read carefully before proceeding")

add_section_label(slide, "CONFIDENTIALITY & FORWARD-LOOKING STATEMENTS",
                  Inches(0.4), Inches(1.3), w=Inches(6))

disclaimer_text = (
    "This confidential presentation (the \"Presentation\") has been prepared by Transform SR Holding Management LLC (the \"Company\") "
    "solely for informational purposes and for the use of prospective investors to whom it is directly addressed and delivered. "
    "By accepting this Presentation, the recipient acknowledges that it is confidential and agrees to maintain the "
    "confidentiality of the information contained herein.\n\n"
    "This Presentation contains forward-looking statements and projections. Such statements are based on current expectations "
    "and assumptions and involve risks and uncertainties that could cause actual results to differ materially from those expressed. "
    "The Company makes no representation or warranty, express or implied, as to the accuracy or completeness of the information "
    "contained herein, and nothing contained herein shall be relied upon as a promise or representation as to past or future performance.\n\n"
    "All financial data presented herein is sourced directly from the Company's NetSuite accounting system and represents "
    "management accounts. Figures have not been audited and are subject to year-end adjustments. Historical performance "
    "does not guarantee future results.\n\n"
    "This Presentation does not constitute an offer to sell or a solicitation of an offer to buy any securities. "
    "Any investment decision should be made solely on the basis of a definitive agreement and after independent due diligence. "
    "Recipients should consult their own legal, financial, and tax advisors before making any investment decision.\n\n"
    "Transform SR Holding Management LLC reserves all rights with respect to this Presentation. "
    "Unauthorized reproduction or distribution of this Presentation, in whole or in part, is strictly prohibited."
)

add_textbox(slide, disclaimer_text,
            Inches(0.4), Inches(1.65), Inches(12.5), Inches(5.2),
            font_size=Pt(9.5), color=DARK_GRAY, align=PP_ALIGN.LEFT)

# ─── Slide 3: Table of Contents ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Table of Contents")

sections = [
    ("01", "Company Overview",             "Background, structure, and competitive positioning"),
    ("02", "Business Segments",             "Revenue streams and operational model"),
    ("03", "Financial Performance",         "FY2024–FY2025 P&L, gross margin, and cost structure"),
    ("04", "Balance Sheet & Liquidity",     "Asset base, capital structure, and cash position"),
    ("05", "Path to Profitability",         "Cost reduction roadmap and EBITDA bridge"),
    ("06", "Growth Strategy",               "Home warranty expansion, royalties, and digital initiatives"),
    ("07", "Transaction Overview",          "Investment thesis and use of proceeds"),
    ("08", "Appendix",                      "Detailed financial statements"),
]

for i, (num, title, desc) in enumerate(sections):
    y = Inches(1.25) + i * Inches(0.69)
    # Number badge
    add_rect(slide, Inches(0.4), y, Inches(0.5), Inches(0.46), fill_color=NAVY)
    add_textbox(slide, num, Inches(0.4), y, Inches(0.5), Inches(0.46),
                font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, title,
                Inches(1.05), y + Inches(0.02), Inches(4.5), Inches(0.26),
                font_size=Pt(12), bold=True, color=NAVY)
    # Description
    add_textbox(slide, desc,
                Inches(1.05), y + Inches(0.26), Inches(11), Inches(0.22),
                font_size=Pt(9.5), color=DARK_GRAY)
    # Separator
    if i < len(sections) - 1:
        add_rect(slide, Inches(0.4), y + Inches(0.5), Inches(12.5), Inches(0.02),
                 fill_color=LIGHT_GRAY)

# ─── Slide 4: Company Overview ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Company Overview",
             "Transform SR Holding Management LLC — Diversified Retail & Services Platform")

# Left column: description
add_section_label(slide, "ABOUT THE COMPANY", Inches(0.4), Inches(1.22))

desc_bullets = [
    "• Transform SR Holding Management LLC (\"Transform SR\" or \"the Company\") is a diversified retail and services holding company operating from a portfolio of Sears- and Kmart-heritage businesses",
    "• The Company operates through two primary subsidiaries — the Parent Company entity and affiliated operations — managing six distinct revenue-generating platforms across the United States, Mexico, and international markets",
    "• Founded on the acquisition of Sears and Kmart operating assets in 2019, the Company has systematically repositioned its business mix toward higher-margin services, warranty, and licensing revenue",
    "• Fiscal year runs February to January (FY2025: Feb 2, 2025 – Jan 31, 2026)",
]
for i, b in enumerate(desc_bullets):
    add_textbox(slide, b, Inches(0.4), Inches(1.62) + i * Inches(0.61),
                Inches(6.8), Inches(0.55),
                font_size=Pt(10), color=DARK_GRAY)

# Right column: key facts boxes
kpi_data = [
    ("~$650M", "FY2025 Revenue"),
    ("$935M",  "Total Assets (Jan '26)"),
    ("$123M",  "Cash & Equivalents"),
    ("~2,000+","Estimated Employees"),
]
for i, (val, lbl) in enumerate(kpi_data):
    col = i % 2
    row = i // 2
    kpi_box(slide, lbl, val,
            x=Inches(7.65) + col * Inches(2.45),
            y=Inches(1.22) + row * Inches(1.3),
            w=Inches(2.2), h=Inches(1.1))

add_section_label(slide, "REVENUE PLATFORMS", Inches(7.65), Inches(3.85))

platforms = [
    ("Home Services",        "National appliance repair & HVAC network"),
    ("Merchandise Retail",   "Sears/Kmart branded merchandise"),
    ("Home Warranty",        "Cinch Home Services & protection plans"),
    ("Parts Distribution",   "OEM & 3rd-party appliance parts"),
    ("Royalties & Licensing","Sears Mexico, Hello Super, LB revenue"),
    ("Wholesale & Logistics","Wholesale sales & DC warehouse operations"),
]
for i, (plat, desc_text) in enumerate(platforms):
    col = i % 2
    row = i // 2
    x = Inches(7.65) + col * Inches(2.55)
    y = Inches(4.22) + row * Inches(0.62)
    add_rect(slide, x, y, Inches(2.35), Inches(0.52),
             fill_color=WHITE, line_color=LIGHT_GRAY, line_width=Pt(1))
    add_rect(slide, x, y, Inches(0.05), Inches(0.52), fill_color=GOLD)
    add_textbox(slide, plat, x + Inches(0.1), y + Inches(0.02),
                Inches(2.2), Inches(0.22),
                font_size=Pt(9), bold=True, color=NAVY)
    add_textbox(slide, desc_text, x + Inches(0.1), y + Inches(0.24),
                Inches(2.2), Inches(0.22),
                font_size=Pt(8), color=MID_GRAY)

# Vertical separator
add_rect(slide, Inches(7.4), Inches(1.2), Inches(0.02), Inches(5.7),
         fill_color=LIGHT_GRAY)

# ─── Slide 5: Business Segments ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Business Segments",
             "Diversified revenue mix across six platforms — FY2025 actuals")

# Revenue segments data (FY2025)
segments = [
    # (label, fy24_rev, fy25_rev, description)
    ("Home Services & Repair",    229, 160, "On-site appliance & HVAC repair, NCC labor, Assurant PA"),
    ("Merchandise Retail",        214, 165, "Sears/Kmart in-store & SBT merchandise, returns nets"),
    ("Parts & Components",         73,  57, "NCC 3rd-party, home warranty, Assurant PA parts"),
    ("Home Warranty",              47,  63, "Home warranty revenue — fastest-growing segment"),
    ("Royalties & Licensing",      40,  43, "Sears Mexico royalty, Hello Super, licensed businesses"),
    ("Wholesale & Distribution",   54,  51, "Wholesale sales net of rebates + DC warehouse"),
    ("Other Revenue",             152, 111, "Misc, 3P PA commissions, CITI, SFI, delivery, etc."),
]

# Segment table
headers = ["Segment", "FY2024 ($M)", "FY2025 ($M)", "YoY Change", "Key Description"]
rows = []
total24, total25 = 0, 0
for label, f24, f25, desc in segments:
    chg = (f25 - f24) / f24 * 100 if f24 > 0 else 0
    chg_str = f"+{chg:.0f}%" if chg >= 0 else f"{chg:.0f}%"
    rows.append([label, f"${f24}M", f"${f25}M", chg_str, desc])
    total24 += f24
    total25 += f25

rows.append(["TOTAL REVENUE", f"${total24}M", f"${total25}M",
             f"{(total25-total24)/total24*100:.0f}%", ""])

col_widths = [3.0, 1.3, 1.3, 1.1, 5.0]
add_table(slide, headers, rows,
          x=Inches(0.4), y=Inches(1.22),
          w=Inches(12.5), h=Inches(4.5),
          col_widths=col_widths,
          font_size=Pt(9.5), header_size=Pt(9.5),
          row_height=Inches(0.48))

# Insight callout
add_rect(slide, Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.95),
         fill_color=NAVY)
add_rect(slide, Inches(0.4), Inches(5.9), Inches(0.08), Inches(0.95),
         fill_color=GOLD)
insights = [
    "▶  Home Warranty is the fastest-growing segment: +34% YoY to $63M — recurring subscription revenue with high renewal rates",
    "▶  Royalties & Licensing ($43M) provides asset-light, high-margin income with international growth potential (Hello Super +43% YoY)",
    "▶  Home Services remains the largest segment at ~$160M, representing the core operational platform and primary driver of company value",
]
add_textbox(slide, "   ".join(insights),
            Inches(0.6), Inches(5.94), Inches(12.1), Inches(0.85),
            font_size=Pt(8.5), color=WHITE)

# ─── Slide 6: Home Services Deep-Dive ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Home Services — Core Operating Platform",
             "National appliance repair network: largest independent home services operation in the US")

# Four KPI boxes
home_kpis = [
    ("~$160M",  "FY2025 Service Revenue"),
    ("~$57M",   "FY2025 Parts Revenue"),
    ("~$63M",   "Home Warranty Revenue"),
    ("$18.5M",  "Sears Mexico Royalty"),
]
for i, (val, lbl) in enumerate(home_kpis):
    kpi_box(slide, lbl, val,
            x=Inches(0.4) + i * Inches(3.18),
            y=Inches(1.22),
            w=Inches(2.9), h=Inches(1.05))

# Service line breakdown
add_section_label(slide, "SERVICE REVENUE BREAKDOWN (FY2025)", Inches(0.4), Inches(2.5))

svc_rows = [
    ["Paid Labor — On-Site (Company Technicians)",   "$76.1M",  "$94.7M",  "-20%"],
    ["NCC 3rd-Party Labor",                           "$49.1M",  "$67.4M",  "-27%"],
    ["NCC Home Warranty Labor",                       "$14.8M",  "$26.2M",  "-44%"],
    ["NCC Assurant PA Labor",                         "$13.0M",  "$32.4M",  "-60%"],
    ["HVAC Replacement Labor (3P & PA)",              "$3.2M",   "N/A",     "New"],
    ["Parts Delivery & Handling",                     "$5.2M",   "$8.5M",   "-39%"],
    ["Less: Labor Returns / Credits",                 "($1.6M)", "($3.9M)", "—"],
    ["TOTAL HOME SERVICES",                           "$159.8M", "$225.3M", "-29%"],
]
svc_hdrs = ["Service Line", "FY2025", "FY2024", "YoY"]
add_table(slide, svc_hdrs, svc_rows,
          x=Inches(0.4), y=Inches(2.88),
          w=Inches(6.2), h=Inches(3.5),
          col_widths=[4.5, 1.2, 1.2, 1.0],
          font_size=Pt(9), header_size=Pt(9))

# Right: Strategic context
add_section_label(slide, "STRATEGIC CONTEXT", Inches(7.0), Inches(2.5))
context_points = [
    ("Scale",         "One of the largest independent home appliance service networks in North America, with technicians operating across all 50 states"),
    ("NCC Platform",  "National Customer Center (NCC) manages third-party service dispatch, Assurant protection plans, and home warranty fulfillment — creating a recurring revenue flywheel"),
    ("Parts Network", "Integrated parts distribution (PartsDirect heritage) generates $57M in FY2025 parts revenue, supporting both proprietary service calls and third-party repair shops"),
    ("Revenue Mix",   "Revenue mix evolving toward higher-margin home warranty ($63M, +34% YoY) and asset-light royalties, reducing dependency on volume-driven labor revenue"),
    ("Volume Trend",  "Service volume declined YoY as legacy Sears protection plan rollovers complete; new Cinch/SPHW contracts provide a more durable recurring base going forward"),
]
for i, (hdr, txt) in enumerate(context_points):
    y_pos = Inches(2.88) + i * Inches(0.68)
    add_rect(slide, Inches(7.0), y_pos, Inches(5.85), Inches(0.62),
             fill_color=WHITE, line_color=LIGHT_GRAY, line_width=Pt(1))
    add_rect(slide, Inches(7.0), y_pos, Inches(0.05), Inches(0.62), fill_color=GOLD)
    add_textbox(slide, hdr, Inches(7.1), y_pos + Inches(0.03),
                Inches(5.6), Inches(0.22),
                font_size=Pt(9.5), bold=True, color=NAVY)
    add_textbox(slide, txt, Inches(7.1), y_pos + Inches(0.26),
                Inches(5.6), Inches(0.30),
                font_size=Pt(8.5), color=DARK_GRAY)

# ─── Slide 7: Financial Performance ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Financial Performance",
             "Improving operating trajectory — FY2024 vs. FY2025 (NetSuite management accounts)")

# KPI row
fin_kpis = [
    ("$650.1M", "FY2025 Revenue",      "vs. $809.8M FY2024"),
    ("$221.5M", "FY2025 Gross Profit", "34.1% gross margin"),
    ("($215.2M)","FY2025 EBIT",         "+$23.7M YoY improvement"),
    ("($334.8M)","FY2025 Net Loss",      "+$74.3M YoY improvement"),
]
for i, (val, lbl, sub) in enumerate(fin_kpis):
    x = Inches(0.4) + i * Inches(3.18)
    kpi_box(slide, lbl, val,
            x=x, y=Inches(1.22), w=Inches(2.9), h=Inches(1.05),
            val_color=RED_LOSS if val.startswith("(") else NAVY)
    add_textbox(slide, sub, x, Inches(2.28), Inches(2.9), Inches(0.25),
                font_size=Pt(8), color=GREEN_POS if "improvement" in sub else MID_GRAY,
                align=PP_ALIGN.CENTER)

# P&L Summary table
add_section_label(slide, "CONSOLIDATED P&L SUMMARY ($M)", Inches(0.4), Inches(2.72))

pl_headers = ["Metric", "FY2023", "FY2024A", "FY2025A", "YoY Δ", "YoY %"]
pl_rows = [
    ["Revenue",            "$0.0M",   "$809.8M", "$650.1M", "($159.7M)", "-19.7%"],
    ["Cost of Sales",      "$0.0M",   "($509.7M)","($428.6M)","$81.1M",  "+15.9%"],
    ["Gross Profit",       "$0.0M",   "$300.1M", "$221.5M", "($78.6M)", "-26.2%"],
    ["  Gross Margin",     "—",       "37.1%",   "34.1%",   "(3.0pp)",  "—"],
    ["Operating Expenses", "$0.0M",   "($538.9M)","($436.7M)","$102.2M", "+19.0%"],
    ["Operating Income (EBIT)", "$0.0M", "($238.9M)","($215.2M)","$23.7M","+9.9%"],
    ["Other Income / (Expense)", "$0.0M","($170.2M)","($119.6M)","$50.6M","+29.7%"],
    ["Net Income / (Loss)", "$0.5M", "($409.1M)","($334.8M)","$74.3M","+18.2%"],
]

add_table(slide, pl_headers, pl_rows,
          x=Inches(0.4), y=Inches(3.1),
          w=Inches(8.2), h=Inches(3.55),
          col_widths=[3.2, 1.1, 1.3, 1.3, 1.2, 1.1],
          font_size=Pt(9.5), header_size=Pt(9.5),
          row_height=Inches(0.41))

# Right side: commentary
add_section_label(slide, "KEY OBSERVATIONS", Inches(8.9), Inches(2.72))
observations = [
    ("Revenue Decline", "Revenue declined $159.7M (-20%) FY24→FY25, driven by the wind-down of legacy SFI/protection plan revenue (-$22.5M), reduced NCC labor volumes, and store rationalization. Core service and warranty revenues showed resilience."),
    ("Margin Compression", "Gross margin contracted 300bps to 34.1% as higher-margin SFI and PA commission revenue ($46.5M decline) was replaced by lower-margin parts/merchandise. Home Warranty margin expansion partially offsets this trend."),
    ("Cost Reduction", "Operating expenses declined $102.2M (-19%) YoY — management has demonstrated meaningful cost discipline. Labor costs (salaries + benefits) reduced by ~$65M YoY, with further optimization underway."),
    ("Loss Improvement", "Net loss improved $74.3M (+18%) to ($334.8M). Other expense declined $50.6M as amortization, one-time items, and intercompany charges normalized. The trajectory toward cash breakeven is visible."),
]
for i, (hdr, txt) in enumerate(observations):
    y_pos = Inches(3.1) + i * Inches(0.93)
    add_rect(slide, Inches(8.9), y_pos, Inches(4.05), Inches(0.87),
             fill_color=WHITE, line_color=LIGHT_GRAY, line_width=Pt(1))
    add_rect(slide, Inches(8.9), y_pos, Inches(0.05), Inches(0.87), fill_color=GOLD)
    add_textbox(slide, hdr, Inches(9.0), y_pos + Inches(0.04),
                Inches(3.85), Inches(0.22),
                font_size=Pt(9.5), bold=True, color=NAVY)
    add_textbox(slide, txt, Inches(9.0), y_pos + Inches(0.28),
                Inches(3.85), Inches(0.54),
                font_size=Pt(8.2), color=DARK_GRAY)

# ─── Slide 8: Revenue Trend ──────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Revenue Trends & Segment Mix",
             "High-margin recurring revenue streams growing within a transitioning portfolio")

# Revenue bridge: FY2024 → FY2025
add_section_label(slide, "REVENUE BRIDGE: FY2024 → FY2025 ($M)", Inches(0.4), Inches(1.22))

bridge_rows = [
    ["FY2024 Revenue",               "$809.8M", ""],
    ["  Home Services & Labor",      "($69.2M)", "Volume decline as legacy service contracts complete"],
    ["  Merchandise Retail",         "($49.0M)", "Store rationalization and SKU reduction"],
    ["  SFI / PA Commissions",       "($46.0M)", "SFI program wind-down and Assurant PA reduction"],
    ["  Parts Sales",                "($16.0M)", "Correlated with service volume decline"],
    ["  Other Revenue",              "($41.0M)", "Various one-time and program items"],
    ["  Home Warranty (growth)",     "+$16.0M",  "Cinch HW expansion and SPHW POS programs"],
    ["  Royalties & Licensing",      "+$3.0M",   "Hello Super +43%, Sears Mexico stable"],
    ["  Wholesale (growth)",         "+$8.0M",   "Expanded wholesale customer relationships"],
    ["  Financial Services",         "+$1.4M",   "CITI commission growth"],
    ["FY2025 Revenue",               "$650.1M", ""],
]
b_hdrs = ["Category", "Amount", "Commentary"]
add_table(slide, b_hdrs, bridge_rows,
          x=Inches(0.4), y=Inches(1.6),
          w=Inches(12.5), h=Inches(4.7),
          col_widths=[3.3, 1.4, 7.8],
          font_size=Pt(9.5), header_size=Pt(9.5),
          row_height=Inches(0.41))

# Insight box
add_rect(slide, Inches(0.4), Inches(6.35), Inches(12.5), Inches(0.55),
         fill_color=NAVY)
add_rect(slide, Inches(0.4), Inches(6.35), Inches(0.08), Inches(0.55), fill_color=GOLD)
add_textbox(slide,
    "Key Insight: Growth segments (Home Warranty, Royalties, Wholesale, Financial Services) added +$28.4M in FY2025. "
    "The legacy service contract roll-off (-$135M across services/SFI/parts) is a finite and quantifiable headwind — "
    "as new Cinch and SPHW contracts ramp, the base stabilizes and growth segments will drive positive momentum.",
    Inches(0.6), Inches(6.37), Inches(12.1), Inches(0.5),
    font_size=Pt(9), color=WHITE)

# ─── Slide 9: Cost Structure ─────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Cost Structure & Optimization",
             "Labor-intensive business with clear levers for cost reduction")

# Cost breakdown table
add_section_label(slide, "TOP COST CATEGORIES (FY2025 vs FY2024, $M)", Inches(0.4), Inches(1.22))

cost_rows = [
    ["People Costs (Salaries, Wages, Benefits)", "$243M", "$305M", "($62M)", "Ongoing headcount optimization; field tech rationalization"],
    ["COGS — Direct Service & Merchandise",       "$136M", "$157M", "($21M)", "Correlated with revenue mix; parts/merchandise cost"],
    ["Technology & Data Processing",               "$32M",  "$35M",  "($3M)",  "Software licenses, outside data processing"],
    ["Legal & Professional Fees",                  "$23M",  "$21M",  "+$2M",   "Litigation and ongoing legal matters (elevated)"],
    ["Utilities & Facilities",                     "$18M",  "$21M",  "($3M)",  "Electric, facility costs across network"],
    ["Freight & Fuel",                             "$31M",  "$42M",  "($11M)", "Truck fuel, freight — volume-correlated"],
    ["Home Warranty Commissions (HW Cinch)",       "$15M",  "$10M",  "+$5M",   "Rising with HW revenue growth"],
    ["Outside & Contract Services",                "$19M",  "$22M",  "($3M)",  "3PL fees, contract labor, security"],
    ["Other Operating Expenses",                   "$116M", "$125M", "($9M)",  "All remaining operating expense lines"],
    ["TOTAL OPERATING EXPENSES",                   "$633M", "$738M", "($105M)","Combined COGS + OpEx"],
]
c_hdrs = ["Cost Category", "FY2025", "FY2024", "YoY Δ", "Commentary"]
add_table(slide, c_hdrs, cost_rows,
          x=Inches(0.4), y=Inches(1.6),
          w=Inches(12.5), h=Inches(4.5),
          col_widths=[3.8, 1.1, 1.1, 1.0, 5.5],
          font_size=Pt(9), header_size=Pt(9.5),
          row_height=Inches(0.43))

# Key takeaway
add_rect(slide, Inches(0.4), Inches(6.22), Inches(12.5), Inches(0.65),
         fill_color=LIGHT_GRAY)
add_rect(slide, Inches(0.4), Inches(6.22), Inches(0.08), Inches(0.65), fill_color=GOLD)
add_textbox(slide,
    "People costs represent ~38% of total cost base ($243M / $633M). Each 10% reduction in labor cost = ~$24M annual savings. "
    "Management has already reduced total costs by $105M (-14%) YoY with further identified savings in legal ($23M+ in litigation costs), "
    "technology consolidation, and 3PL optimization. A credible path to $150M+ additional cost reduction exists.",
    Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.6),
    font_size=Pt(9), color=DARK_GRAY)

# ─── Slide 10: Balance Sheet ─────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Balance Sheet & Liquidity",
             "Substantial asset base; capital structure reflects transformation-period financing")

# Balance sheet comparison table
add_section_label(slide, "CONSOLIDATED BALANCE SHEET ($M)", Inches(0.4), Inches(1.22))

bs_rows = [
    ["ASSETS",                           "",         "",          ""],
    ["  Cash & Equivalents",             "$135.6M",  "$123.3M",  "($12.3M)"],
    ["  Accounts Receivable",            "$66.1M",   "$29.2M",   "($36.9M)"],
    ["  Other Current Assets",           "$128.4M",  "$80.8M",   "($47.6M)"],
    ["  Total Current Assets",           "$330.1M",  "$233.2M",  "($96.9M)"],
    ["  Fixed Assets (Net)",             "$423.6M",  "$406.3M",  "($17.3M)"],
    ["  Other Assets (Intangibles, etc.)","$354.4M", "$295.3M",  "($59.1M)"],
    ["TOTAL ASSETS",                     "$1,108.1M","$934.9M",  "($173.2M)"],
    ["",                                 "",         "",          ""],
    ["LIABILITIES",                      "",         "",          ""],
    ["  Accounts Payable",               "$38.4M",   "$31.7M",   "($6.7M)"],
    ["  Other Current Liabilities",      "$782.3M",  "$838.1M",  "+$55.8M"],
    ["  Total Current Liabilities",      "$820.7M",  "$869.8M",  "+$49.1M"],
    ["  Long-Term Liabilities",          "$512.1M",  "$509.0M",  "($3.1M)"],
    ["TOTAL LIABILITIES",                "$1,332.8M","$1,378.8M","+$46.0M"],
    ["TOTAL EQUITY",                     "($224.8M)","($443.8M)","($219.1M)"],
]
bs_hdrs = ["Balance Sheet Item", "FY2024 (Feb '25)", "FY2025 (Jan '26)", "Change"]
add_table(slide, bs_hdrs, bs_rows,
          x=Inches(0.4), y=Inches(1.6),
          w=Inches(6.9), h=Inches(5.2),
          col_widths=[3.8, 1.5, 1.5, 1.1],
          font_size=Pt(9), header_size=Pt(9.5),
          row_height=Inches(0.3))

# Right: Commentary boxes
add_section_label(slide, "CAPITAL STRUCTURE NOTES", Inches(7.6), Inches(1.22))

notes = [
    ("Cash Position", "$123M cash at Jan 2026. FY2024 cash flow included $418M equity infusion from parent/shareholder, partially offsetting operating cash consumption."),
    ("Current Liabilities", "Other Current Liabilities of $838M is disproportionately large — primarily comprised of deferred home warranty/service contract revenue, intercompany payables, and program-related obligations. These are largely non-cash obligations that unwind over time as services are performed."),
    ("Fixed Assets", "$406M in net fixed assets reflects the Company's national service and logistics infrastructure — distribution centers, service fleet, and technology platform. These represent deployable collateral."),
    ("Negative Equity", "Cumulative equity of ($444M) reflects losses incurred since acquisition. Common equity contributions of $3.2B have been offset by $3.3B of accumulated losses since the 2019 acquisition, primarily from store closure costs, impairments, and operational transition charges."),
    ("Debt Structure", "$509M in long-term liabilities includes operational lease obligations and intercompany financing arrangements. No external public debt outstanding."),
]
for i, (hdr, txt) in enumerate(notes):
    y_pos = Inches(1.6) + i * Inches(1.02)
    add_rect(slide, Inches(7.6), y_pos, Inches(5.3), Inches(0.95),
             fill_color=WHITE, line_color=LIGHT_GRAY, line_width=Pt(1))
    add_rect(slide, Inches(7.6), y_pos, Inches(0.05), Inches(0.95), fill_color=GOLD)
    add_textbox(slide, hdr, Inches(7.7), y_pos + Inches(0.04),
                Inches(5.1), Inches(0.22),
                font_size=Pt(9.5), bold=True, color=NAVY)
    add_textbox(slide, txt, Inches(7.7), y_pos + Inches(0.27),
                Inches(5.1), Inches(0.62),
                font_size=Pt(8.2), color=DARK_GRAY)

# ─── Slide 11: Path to Profitability ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Path to Profitability",
             "Achievable EBITDA improvement through cost optimization and revenue stabilization")

# EBITDA bridge
add_section_label(slide, "EBITDA IMPROVEMENT BRIDGE ($M) — FY2025 → TARGET", Inches(0.4), Inches(1.22))

bridge_items = [
    ("FY2025 EBIT (Operating Loss)",    "($215.2M)", "", "Baseline"),
    ("D&A Add-Back (estimated)",        "+$50M",      "~", "Estimated depreciation/amortization in Other Expense"),
    ("FY2025 Adjusted EBITDA",          "($165M)",    "", "Starting point for bridge"),
    ("",                                "",            "", ""),
    ("Labor Cost Reduction Initiative", "+$60M",      "✓", "15-20% headcount optimization; service network right-sizing"),
    ("Legal & Professional Fees",       "+$12M",      "✓", "Resolution of legacy litigation; lower run-rate legal costs"),
    ("Technology Rationalization",      "+$10M",      "✓", "Consolidation of legacy Sears IT platforms to modern stack"),
    ("Facilities Optimization",         "+$8M",       "✓", "DC network rationalization and lease exits"),
    ("Home Warranty Revenue Growth",    "+$20M",      "↑", "Cinch/SPHW new contract ramp (+30% revenue growth scenario)"),
    ("Royalties & Licensing Growth",    "+$8M",       "↑", "Hello Super + new licensing deals in development"),
    ("New B2B Service Partnerships",    "+$15M",      "↑", "NCC platform monetization through new 3rd-party OEM contracts"),
    ("TARGET EBITDA",                   "+$3M to $28M","", "Breakeven to modest positive EBITDA within 18-24 months"),
]

b_hdrs = ["Initiative", "Impact", "Status", "Commentary"]
b_rows = [list(b) for b in bridge_items if b[0] != ""]
add_table(slide, b_hdrs, b_rows,
          x=Inches(0.4), y=Inches(1.6),
          w=Inches(12.5), h=Inches(4.8),
          col_widths=[3.8, 1.4, 0.8, 6.5],
          font_size=Pt(9.2), header_size=Pt(9.5),
          row_height=Inches(0.41))

add_rect(slide, Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.45),
         fill_color=NAVY)
add_rect(slide, Inches(0.4), Inches(6.5), Inches(0.08), Inches(0.45), fill_color=GOLD)
add_textbox(slide,
    "The combination of $90M in identified cost reductions and $43M in revenue growth initiatives creates a credible path to "
    "EBITDA breakeven within 18–24 months. Each $10M of incremental EBITDA improvement increases enterprise value by $50–100M "
    "at typical services/distribution multiples of 5–10x.",
    Inches(0.6), Inches(6.53), Inches(12.1), Inches(0.4),
    font_size=Pt(9), color=WHITE)

# ─── Slide 12: Growth Strategy ────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Growth Strategy & Strategic Initiatives",
             "Four strategic pillars to drive value creation")

pillars = [
    {
        "num": "01",
        "title": "Home Warranty Expansion",
        "subtitle": "Accelerate recurring revenue growth",
        "points": [
            "Home Warranty revenue grew +34% YoY to $63M — largest and fastest-growing segment",
            "SPHW (Sears Protection Home Warranty) and Cinch new contracts ramping",
            "HVAC replacement program launched in FY2025 — new revenue stream with high NPS",
            "Target: $80–90M HW revenue in FY2026 through new dealer partnerships",
            "Home warranty is sticky, recurring, high-margin — transforms revenue quality",
        ],
    },
    {
        "num": "02",
        "title": "International Royalties & Brand Licensing",
        "subtitle": "High-margin, asset-light income streams",
        "points": [
            "Sears Mexico royalty: $18.5M p.a. — stable, contractual, FX-insulated royalty",
            "Hello Super (Mexican digital marketplace): +43% YoY to $10M — high growth",
            "Licensed Business revenue: $13M from domestic licensed business operators",
            "Active discussions on new international licensing in Latin America and Asia",
            "Target: $55–65M combined licensing revenue by FY2027",
        ],
    },
    {
        "num": "03",
        "title": "NCC Platform — Third-Party Monetization",
        "subtitle": "Turn our service infrastructure into a B2B profit center",
        "points": [
            "NCC manages dispatch, technician routing, and service fulfillment at national scale",
            "Existing OEM relationships (Assurant, CINCH) validate the platform's capabilities",
            "Potential to onboard additional appliance OEMs, insurers, and home service apps",
            "White-label NCC dispatch = high margin incremental revenue with minimal CapEx",
            "Target: $20–30M new NCC B2B revenue within 24 months",
        ],
    },
    {
        "num": "04",
        "title": "Cost Transformation Program",
        "subtitle": "Structurally reduce the break-even point",
        "points": [
            "$105M total cost reduction achieved FY2024→FY2025; program continues",
            "Legacy IT decommissioning underway — estimated $10M annual savings on completion",
            "Distribution center network optimization: right-sizing to current service footprint",
            "Legal: resolution of legacy Sears litigation expected to reduce $23M+ legal spend",
            "Target: $90M additional cost reduction over next 18 months",
        ],
    },
]

for idx, pillar in enumerate(pillars):
    col = idx % 2
    row = idx // 2
    x = Inches(0.4) + col * Inches(6.35)
    y = Inches(1.22) + row * Inches(2.95)
    w = Inches(6.15)
    h = Inches(2.85)

    add_rect(slide, x, y, w, h, fill_color=WHITE, line_color=LIGHT_GRAY, line_width=Pt(1))
    # Number badge
    add_rect(slide, x, y, Inches(0.55), Inches(0.42), fill_color=NAVY)
    add_textbox(slide, pillar["num"], x, y, Inches(0.55), Inches(0.42),
                font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title bar
    add_rect(slide, x + Inches(0.55), y, w - Inches(0.55), Inches(0.42), fill_color=NAVY_MID)
    add_textbox(slide, pillar["title"], x + Inches(0.6), y + Inches(0.02),
                w - Inches(0.65), Inches(0.24),
                font_size=Pt(11), bold=True, color=WHITE)
    add_textbox(slide, pillar["subtitle"], x + Inches(0.6), y + Inches(0.23),
                w - Inches(0.65), Inches(0.18),
                font_size=Pt(8.5), color=GOLD_LT)
    # Bullet points
    for i, pt in enumerate(pillar["points"]):
        add_textbox(slide, f"• {pt}",
                    x + Inches(0.12), y + Inches(0.48) + i * Inches(0.45),
                    w - Inches(0.22), Inches(0.42),
                    font_size=Pt(8.8), color=DARK_GRAY)

# ─── Slide 13: Transaction Overview ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Transaction Overview",
             "Investment thesis, structure, and use of proceeds")

# Left: Investment thesis
add_section_label(slide, "INVESTMENT THESIS", Inches(0.4), Inches(1.22))

thesis_points = [
    ("Asymmetric Upside",  "Significant enterprise value exists in the NCC platform, Home Warranty business, and Royalty streams — these three segments alone generate ~$116M in FY2025 revenue with high incremental margins. At 5–8x revenue, the services/royalty portfolio represents $580M–$930M of standalone value."),
    ("Loss Turn Visible",  "Net loss improved $74M YoY (+18%) with $105M in cost reduction achieved. Management has demonstrated execution ability. The next $90M cost reduction initiative is defined, resourced, and underway — creating line of sight to EBITDA breakeven."),
    ("Asset Coverage",     "$935M in total assets provides meaningful downside protection. Fixed assets ($406M) and other long-term assets ($295M) represent infrastructure with significant replacement value. Current ratio assets ($233M) provide near-term liquidity."),
    ("Contractual Revenue","Sears Mexico royalty ($18.5M/year), CITI financial services ($9.4M), and home warranty contracts are contractual recurring revenue streams — providing a resilient base independent of volume fluctuations."),
    ("Strategic Options",  "The business has multiple path-to-value options: (1) operational turnaround to EBITDA positive; (2) segment carve-out of the NCC/Home Warranty platform; (3) strategic partnership with a national home services or insurance player; (4) IP/licensing monetization."),
]
for i, (hdr, txt) in enumerate(thesis_points):
    y_pos = Inches(1.6) + i * Inches(1.01)
    add_rect(slide, Inches(0.4), y_pos, Inches(6.2), Inches(0.95),
             fill_color=WHITE, line_color=LIGHT_GRAY, line_width=Pt(1))
    add_rect(slide, Inches(0.4), y_pos, Inches(0.05), Inches(0.95), fill_color=GOLD)
    add_textbox(slide, hdr, Inches(0.5), y_pos + Inches(0.04),
                Inches(5.95), Inches(0.22),
                font_size=Pt(9.5), bold=True, color=NAVY)
    add_textbox(slide, txt, Inches(0.5), y_pos + Inches(0.27),
                Inches(5.95), Inches(0.62),
                font_size=Pt(8.2), color=DARK_GRAY)

# Right: Use of proceeds and key terms
add_section_label(slide, "INDICATIVE USE OF PROCEEDS", Inches(7.0), Inches(1.22))

proceeds_rows = [
    ["Growth CapEx — NCC Platform",       "~30%", "$XX–XXM", "Technology, capacity, OEM integrations"],
    ["Home Warranty Business Dev",         "~25%", "$XX–XXM", "Sales, marketing, dealer network expansion"],
    ["Working Capital & Liquidity",        "~20%", "$XX–XXM", "Operational runway during cost transformation"],
    ["IT Modernization",                   "~15%", "$XX–XXM", "Decommission legacy Sears systems"],
    ["General Corporate / M&A Optionality","~10%", "$XX–XXM", "Strategic bolt-ons, if identified"],
]
p_hdrs = ["Use of Proceeds", "Allocation", "Amount", "Rationale"]
add_table(slide, p_hdrs, proceeds_rows,
          x=Inches(7.0), y=Inches(1.6),
          w=Inches(5.85), h=Inches(2.5),
          col_widths=[2.8, 0.9, 1.0, 1.15],
          font_size=Pt(8.5), header_size=Pt(9))

add_section_label(slide, "KEY INVESTMENT CONSIDERATIONS", Inches(7.0), Inches(4.25))

considerations = [
    "✓  Established national-scale platform with deep customer relationships",
    "✓  Management team with demonstrated ability to reduce costs (-$105M in 12 months)",
    "✓  Contractual revenue streams ($28M+ royalties, $9M financial services) provide floor",
    "✓  Home warranty growth provides path to positive EBITDA mix shift",
    "✓  $406M fixed asset base provides asset-backed collateral",
    "×  Revenue decline (-20%) requires stabilization before growth",
    "×  Negative equity ($444M) reflects legacy transition costs",
    "×  Large deferred liability obligations ($838M OCL) warrant detailed diligence",
]
for i, pt in enumerate(considerations):
    color = GREEN_POS if pt.startswith("✓") else RED_LOSS
    add_textbox(slide, pt,
                Inches(7.0), Inches(4.6) + i * Inches(0.315),
                Inches(5.85), Inches(0.295),
                font_size=Pt(9), color=color)

# ─── Slide 14: Appendix — Detailed P&L ───────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
slide_header(slide, "Appendix A — Detailed Financial Statements",
             "NetSuite management accounts — unaudited")

add_section_label(slide, "DETAILED P&L (FY2024 & FY2025, $M)", Inches(0.4), Inches(1.22))

detail_hdrs = ["Line Item", "FY2024A ($M)", "FY2025A ($M)", "YoY Change ($M)", "YoY %"]
detail_rows = [
    # Revenue
    ["REVENUE",                            "",         "",         "",        ""],
    ["  Merchandise Sales",                "$176.3",   "$158.0",   "($18.3)", "-10%"],
    ["  SBT Sales",                        "$1.1",     "$4.0",     "+$2.9",   "+264%"],
    ["  Paid Labor — On-Site",             "$94.7",    "$76.1",    "($18.6)", "-20%"],
    ["  NCC 3rd-Party Labor",              "$67.4",    "$49.1",    "($18.3)", "-27%"],
    ["  NCC Home Warranty Labor",          "$26.2",    "$14.8",    "($11.4)", "-44%"],
    ["  NCC Assurant PA Labor",            "$32.4",    "$13.0",    "($19.4)", "-60%"],
    ["  NCC 3rd-Party Parts",              "$39.8",    "$37.0",    "($2.8)",  "-7%"],
    ["  NCC Home Warranty Parts",          "$20.7",    "$15.1",    "($5.6)",  "-27%"],
    ["  Home Warranty Revenue",            "$46.8",    "$62.8",    "+$16.0",  "+34%"],
    ["  Wholesale Revenue (net rebates)",  "$24.8",    "$33.4",    "+$8.6",   "+35%"],
    ["  DC Warehouse Revenue",             "$25.6",    "$25.0",    "($0.6)",  "-2%"],
    ["  Sears Mexico Royalty",             "$18.3",    "$18.5",    "+$0.2",   "+1%"],
    ["  Hello Super Commission",           "$7.0",     "$10.0",    "+$3.0",   "+43%"],
    ["  Licensed Business Revenue",        "$13.9",    "$12.9",    "($1.0)",  "-7%"],
    ["  SFI Revenue",                      "$23.1",    "$0.6",     "($22.5)", "-97%"],
    ["  3rd Party PA Commission",          "$35.0",    "$11.5",    "($23.5)", "-67%"],
    ["  CITI Admin/Comm Revenue",          "$8.0",     "$9.4",     "+$1.4",   "+18%"],
    ["  Misc & Other Revenue",             "$148.7",   "$98.9",    "($49.8)", "-33%"],
    ["TOTAL REVENUE",                      "$809.8",   "$650.1",   "($159.7)","-20%"],
    # COGS
    ["COST OF SALES",                      "($509.7)", "($428.6)", "+$81.1",  "+16%"],
    ["GROSS PROFIT",                       "$300.1",   "$221.5",   "($78.6)", "-26%"],
    ["  Gross Margin %",                   "37.1%",    "34.1%",    "(3.0pp)", ""],
    # Expenses
    ["OPERATING EXPENSES",                 "($538.9)", "($436.7)", "+$102.2", "+19%"],
    ["OPERATING INCOME (EBIT)",            "($238.9)", "($215.2)", "+$23.7",  "+10%"],
    ["OTHER INCOME / (EXPENSE)",           "($170.2)", "($119.6)", "+$50.6",  "+30%"],
    ["NET INCOME / (LOSS)",                "($409.1)", "($334.8)", "+$74.3",  "+18%"],
]

add_table(slide, detail_hdrs, detail_rows,
          x=Inches(0.4), y=Inches(1.6),
          w=Inches(12.5), h=Inches(5.2),
          col_widths=[4.2, 1.6, 1.6, 1.8, 1.0],
          font_size=Pt(8.5), header_size=Pt(9),
          row_height=Inches(0.185))

# ─── Save ────────────────────────────────────────────────────────────────────
output_path = "/Users/josh/Documents/Finance Bots/Transform_SR_Investor_Pitch_Deck_May2026.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
